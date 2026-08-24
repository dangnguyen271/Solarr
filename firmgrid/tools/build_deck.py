"""Generate the FirmGrid presentation deck + TV poster (16:9 PPTX).

Run:  python build_deck.py
Structure follows the 6-part hackathon pitch framework:
intro / problem / impact / product / demo / why-now / wrap-up.
Minimal text — key points and figures only; the full script lives in
each slide's speaker notes.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "submission"
OUT.mkdir(exist_ok=True)

W, H = Inches(13.333), Inches(7.5)

NAVY = RGBColor(0x0B, 0x12, 0x20)
PANEL = RGBColor(0x10, 0x1A, 0x30)
PANEL_EDGE = RGBColor(0x24, 0x3B, 0x5E)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x38, 0xBD, 0xF8)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
INK = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
FONT = "Avenir Next"


# --------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------- #
def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Pt(6), W, Pt(6))
    strip.fill.solid()
    strip.fill.fore_color.rgb = GREEN
    strip.line.fill.background()
    strip.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, runs, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
        space_after=6, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, color, bold)]
        for text, c, b in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.color.rgb = c
            r.font.name = FONT
    return box


def kicker(slide, label, timing=None):
    t = f"FIRMGRID  ·  {label}" + (f"  ·  {timing}" if timing else "")
    txt(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
        [[(t, MUTED, True)]], size=12)


def headline(slide, text_, size=32, y=Inches(0.72), color=INK, w=Inches(12.1)):
    txt(slide, Inches(0.6), y, w, Inches(1.2), [[(text_, color, True)]], size=size)


def panel(slide, x, y, w, h, edge=PANEL_EDGE):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.06
    p.fill.solid()
    p.fill.fore_color.rgb = PANEL
    p.line.color.rgb = edge
    p.line.width = Pt(1.2)
    p.shadow.inherit = False
    return p


def stat_tile(slide, x, y, w, h, big, label, accent=GREEN, big_size=32):
    panel(slide, x, y, w, h)
    txt(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.85),
        [[(big, accent, True)]], size=big_size)
    txt(slide, x + Inches(0.18), y + h - Inches(0.92), w - Inches(0.36), Inches(0.85),
        [[(label, MUTED, False)]], size=12.5, space_after=0)


def bullets(slide, x, y, w, h, items, size=16, gap=9):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("▸ ", GREEN, True), (lead, INK, True), (rest, INK, False)])
        else:
            paras.append([("▸ ", GREEN, True), (it, INK, False)])
    txt(slide, x, y, w, h, paras, size=size, space_after=gap)


def notes(slide, text_):
    slide.notes_slide.notes_text_frame.text = text_


def flow_box(slide, x, y, w, h, title, sub, accent=BLUE):
    panel(slide, x, y, w, h, edge=accent)
    txt(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.4),
        [[(title, accent, True)]], size=13, space_after=0)
    txt(slide, x + Inches(0.1), y + Inches(0.48), w - Inches(0.2), h - Inches(0.56),
        [[(sub, MUTED, False)]], size=10.5, space_after=0)


def arrow(slide, x, y, w=Inches(0.3)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.26))
    a.fill.solid()
    a.fill.fore_color.rgb = GREEN
    a.line.fill.background()
    a.shadow.inherit = False


# --------------------------------------------------------------------- #
# deck
# --------------------------------------------------------------------- #
prs = Presentation()
prs.slide_width, prs.slide_height = W, H

# ---- 1 · INTRO (30s) ------------------------------------------------- #
s = add_slide(prs)
txt(s, Inches(0.6), Inches(0.9), Inches(12.1), Inches(1.0),
    [[("⚡ FirmGrid", GREEN, True)]], size=52)
txt(s, Inches(0.6), Inches(2.05), Inches(12.1), Inches(1.2),
    [[("Vietnam throws away clean power at noon —", INK, True)],
     [("and charges its clean vehicles on coal at night.", INK, True)]], size=26)
txt(s, Inches(0.6), Inches(3.35), Inches(12.1), Inches(0.6),
    [[("Sun-to-Wheels today. Sun-to-Servers next.", GREEN, True)]], size=18)
mw = Inches(2.92)
for i, (name, role) in enumerate([
    ("[Member 1]", "Team lead · ML & optimisation"),
    ("[Member 2]", "Power systems & digital twin"),
    ("[Member 3]", "Backend & settlement"),
    ("[Member 4]", "Product, UX & policy"),
]):
    x = Inches(0.6) + i * (mw + Inches(0.15))
    panel(s, x, Inches(4.3), mw, Inches(1.8), edge=BLUE if i % 2 else GREEN)
    txt(s, x + Inches(0.15), Inches(4.5), mw - Inches(0.3), Inches(0.6),
        [[(name, INK, True)]], size=16)
    txt(s, x + Inches(0.15), Inches(5.15), mw - Inches(0.3), Inches(0.8),
        [[(role, MUTED, False)]], size=12)
txt(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
    [[("[TEAM_NAME]  ·  Asian Hackathon for Green Future 2026  ·  Track 1", MUTED, False)]],
    size=13)
notes(s, "INTRO (30s). Welcome the judges, introduce each member by name and what they "
         "built — one sentence each. Then the hook: every sunny morning Vietnam curtails "
         "rooftop solar; every evening its new electric motorbikes charge on coal. "
         "FirmGrid closes both loops — and everything you'll see runs live on this laptop.")

# ---- 2 · PROBLEM (30s, root cause folded in) -------------------------- #
s = add_slide(prs)
kicker(s, "PROBLEM", "30s")
headline(s, "Wasted at noon. Coal at night. And 2026 makes it worse.")
y, w3 = Inches(1.95), Inches(3.93)
stat_tile(s, Inches(0.6), y, w3, Inches(1.85), "~19 GW",
          "of solar already built — a year's power for ~10 million homes — "
          "yet curtailed feeder-wide on sunny mornings", AMBER)
stat_tile(s, Inches(4.7), y, w3, Inches(1.85), "450,000",
          "petrol bikes replaced under Hanoi's LEZ — 1 in 15 of the city's motorbikes — "
          "set to charge on evening coal", GREEN)
stat_tile(s, Inches(8.8), y, w3, Inches(1.85), "2× by 2030",
          "data-centre demand: 735 → ~1,500 MW — the draw of ~4 million homes — "
          "mandated ≥50% green", BLUE)
panel(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(2.1), edge=RED)
txt(s, Inches(0.85), Inches(4.45), Inches(11.6), Inches(0.5),
    [[("ROOT CAUSE — curtailment is blunt because it is blind", RED, True)]], size=17)
bullets(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.3), [
    ("Operators can't see which home could keep exporting ", "→ whole feeders are cut."),
    ("26 Jun 2026: surplus-sale cap raised 20% → 50% ",
     "— the ceiling doubled, and nobody has a fair way to allocate it."),
], size=15, gap=8)
notes(s, "PROBLEM (30s). Three figures: 19 GW of solar that gets cut feeder-wide; 450k "
         "motorbikes electrifying into the evening peak; data-centre demand doubling with a "
         "green mandate. Root cause in one line: the operator is blind, so the cut is blunt "
         "— it's an intelligence gap, not a hardware gap. And the surplus cap just doubled "
         "a week before this pitch: more pressure, still no allocator.")

# ---- 3 · IMPACT (30s) -------------------------------------------------- #
s = add_slide(prs)
kicker(s, "IMPACT", "30s")
headline(s, "What one intelligence layer is worth")
y = Inches(2.0)
w4 = Inches(2.9)
stat_tile(s, Inches(0.6), y, w4, Inches(2.0), "15–19 GWh/yr",
          "rescued across Hanoi — a year's electricity for ~6,000 homes "
          "(10–13 kt CO₂)", GREEN, 24)
stat_tile(s, Inches(3.67), y, w4, Inches(2.0), "16–27 kt CO₂/yr",
          "avoided by charging e-bikes on sunshine — like planting a million "
          "trees, every year", BLUE, 24)
stat_tile(s, Inches(6.74), y, w4, Inches(2.0), "0.6–1.0M ₫/yr",
          "back in each solar family's pocket — a month's power bill, from "
          "energy thrown away today", AMBER, 24)
stat_tile(s, Inches(9.81), y, w4, Inches(2.0), "$85 vs $92",
          "per MWh: clean 24/7 power ~8% cheaper than the grid tariff data "
          "centres pay", GREEN, 24)
bullets(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(1.6), [
    ("EVN: ", "transformer reinforcement deferred — hundreds of millions ₫ per site."),
    ("National scenario 2030: ", "300–380 GWh/yr ≈ 200–260 kt CO₂/yr."),
    ("Every assumption is a slider ", "— judges can recompute all of this live."),
], size=15)
notes(s, "IMPACT (30s). Supply side: 15-19 GWh a year recovered across Hanoi's constrained "
         "transformers — 10-13 kt of CO2. Mobility side: another 16-27 kt by steering the "
         "LEZ charging wave into the solar window — one kilogram per swap, printed on the "
         "rider's receipt. Households earn real money from energy that's discarded today; "
         "EVN defers copper. And at Tier 2, the clean block undercuts the grid tariff.")

# ---- 4 · PRODUCT (1 min) ----------------------------------------------- #
s = add_slide(prs)
kicker(s, "PRODUCT", "1 min")
headline(s, "The intelligence layer that makes the single-buyer grid smart",
         color=GREEN, size=28)
txt(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.5),
    [[("USP: legal today (Decrees 58 & 57/2025, VND rails) · zero new household hardware · "
       "explainable to a grandmother", AMBER, True)]], size=15)
fw = Inches(3.93)
feats = [
    ("GridMind", "PREDICTS safe export headroom per transformer, per 15 minutes — "
     "F1 0.86 on real Hanoi weather", BLUE),
    ("HeadRoom Auction", "ALLOCATES it fairly — minimal curtailment, ≤90% of forecast, "
     "priority credits for the declined", GREEN),
    ("FlexMatch", "STEERS swap stations & depots to charge on sunshine — "
     "Sun-to-Wheels", AMBER),
]
for i, (t4, sub4, c4) in enumerate(feats):
    x = Inches(0.6) + i * (fw + Inches(0.17))
    panel(s, x, Inches(2.35), fw, Inches(2.5), edge=c4)
    txt(s, x + Inches(0.18), Inches(2.55), fw - Inches(0.36), Inches(0.5),
        [[(t4, c4, True)]], size=19)
    txt(s, x + Inches(0.18), Inches(3.15), fw - Inches(0.36), Inches(1.6),
        [[(sub4, INK, False)]], size=13.5)
panel(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(1.35), edge=BLUE)
txt(s, Inches(0.85), Inches(5.32), Inches(11.6), Inches(1.1),
    [[("Same engine, Tier 2 — Sun-to-Servers: ", BLUE, True),
      ("aggregate orchestrated feeders + storage into firm 24/7 clean blocks for data "
       "centres via DPPA. The bridge to nuclear-era baseload (2030–35) — we win either way.",
       INK, False)]], size=14.5)
notes(s, "PRODUCT (60s). One sentence: FirmGrid doesn't fight the single-buyer system, it "
         "makes it intelligent. Three features only: GridMind predicts, the auction "
         "allocates fairly, FlexMatch steers mobility charging into the sun. Settlement is "
         "verified, VND-only, hash-chain audited — no crypto. And the same engine scales to "
         "Tier 2: firm hour-matched blocks for data centres. If nuclear lands on time we "
         "built the flexible grid it needs; if it slips, we're why the lights stayed clean.")

# ---- 5 · DEMO (2 min, mechanism merged in) ------------------------------ #
s = add_slide(prs)
kicker(s, "LIVE DEMO — WORKING PROTOTYPE", "2 min")
headline(s, "A real Hanoi sunny day (26 May 2026), twice", size=28)
# mini flow strip
bw, gap_ = Inches(2.72), Inches(0.36)
steps5 = [("SENSE + FORECAST", BLUE), ("BID (capped at physics)", GREEN),
          ("CLEAR (fair MILP)", AMBER), ("SETTLE (VND + ledger)", GREEN)]
for i, (t5, c5) in enumerate(steps5):
    x = Inches(0.6) + i * (bw + gap_)
    p = panel(s, x, Inches(1.7), bw, Inches(0.62), edge=c5)
    txt(s, x, Inches(1.78), bw, Inches(0.45), [[(t5, c5, True)]],
        size=13, align=PP_ALIGN.CENTER)
    if i:
        arrow(s, x - gap_ + Inches(0.03), Inches(1.88))
y = Inches(2.7)
w4 = Inches(2.9)
stat_tile(s, Inches(0.6), y, w4, Inches(1.8), "703 → 20",
          "kWh wasted, before vs with FirmGrid — 97% rescued: a day's power "
          "for ~100 homes", GREEN, 28)
stat_tile(s, Inches(3.67), y, w4, Inches(1.8), "17 → 0",
          "times the transformer hit its red line — zero, while selling MORE "
          "solar", BLUE, 28)
stat_tile(s, Inches(6.74), y, w4, Inches(1.8), "791,000 ₫",
          "paid to local solar families in one sunny day — for energy that "
          "today earns nothing", AMBER, 28)
stat_tile(s, Inches(9.81), y, w4, Inches(1.8), "F1 = 0.86",
          "on real Hanoi weather — catches ~6 of every 7 overloads an hour "
          "ahead", GREEN, 28)
bullets(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.7), [
    ("Safer AND less wasteful: ", "97% of curtailed energy recovered with zero limit "
     "breaches — the safety envelope holds even through a judge-made storm."),
    ("Trust built in: ", "fair queue (no home waits more than 1 window) · inflated bids "
     "blocked before money moves · every payment hash-chain audited."),
    ("Honest by design: ", "real Hanoi weather, held-out metrics on screen, every decision "
     "explained in one sentence."),
], size=15)
notes(s, "DEMO (120s — switch to the app, keep this slide as backup). Choreography: "
         "1) Neighbourhood Live map, scrub to noon — baseline side turns red (feeder-wide "
         "cut), FirmGrid side stays green, stations visibly swell. 2) Tab 1 counters: 703 "
         "kWh wasted becomes 20 — 97% recovered, zero breaches, 791k dong paid. 3) Hand a "
         "judge the mouse: storm slider, then the fraud bid — Sentinel blocks it and seals "
         "the event in the hash-chained ledger. Close with: the weather is real — twelve "
         "months of Hanoi data ending last week. The pipeline strip above is the mechanism: "
         "sense, bid, clear, settle — every 15 minutes, never above 90% of forecast "
         "headroom, human override on top.")

# ---- 6 · WHY NOW (timely) ----------------------------------------------- #
s = add_slide(prs)
kicker(s, "WHY NOW")
headline(s, "Four clocks started ticking this quarter")
bullets(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(3.9), [
    ("1 Jul 2026 — Hanoi LEZ is live: ", "450,000 petrol bikes to replace; charging habits "
     "crystallise in the next 24 months — retrofitting coordination later is far harder."),
    ("26 Jun 2026 — surplus-sale cap 20% → 50%: ", "sellable rooftop energy doubled "
     "overnight; the same constrained transformers must now allocate it."),
    ("Jan 2026 — data-localisation law: ", "domestic data-centre build-out accelerates "
     "(735 MW → ~1,500 MW by 2030), mandated ≥50% green — buyers for firm clean blocks."),
    ("2030–2035 — earliest nuclear (Ninh Thuận 1): ", "the firm-power gap until then is "
     "ours to fill; PDP8's 10–16 GW of storage takes years and billions — software ships "
     "in months."),
], size=16, gap=14)
txt(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7),
    [[("Every quarter of delay = curtailed GWh gone forever and charging habits locked "
       "onto coal.", AMBER, True)]], size=15)
notes(s, "WHY NOW (20s). Four dates: the LEZ went live this week; the surplus cap doubled "
         "last week; the localisation law is pulling data centres onshore with a green "
         "mandate; and nothing firm arrives before 2030-35. The window is now — and every "
         "sunny day that passes, the wasted energy is gone forever.")

# ---- 7 · WRAP UP (30s) --------------------------------------------------- #
s = add_slide(prs)
kicker(s, "WRAP UP", "30s")
y = Inches(1.7)
for i, (lead, rest, c7) in enumerate([
    ("IMPACT", "15–19 GWh/yr + 26–40 kt CO₂/yr at Hanoi scale — recomputable live, "
     "every assumption on a slider.", GREEN),
    ("MARKET FIT", "Regulation-native: EVN stays sole buyer, stations join via DPPA, "
     "VND-only settlement. Ships without new law — pilot-ready with EVNHANOI Q3 2026, "
     "first firm DPPA block H2 2027.", BLUE),
    ("USP", "The only platform that predicts congestion before it happens, allocates "
     "headroom fairly, and turns mobility batteries into the grid's sponge.", AMBER),
]):
    panel(s, Inches(0.6), y, Inches(12.1), Inches(1.28), edge=c7)
    txt(s, Inches(0.85), y + Inches(0.14), Inches(2.2), Inches(0.5),
        [[(lead, c7, True)]], size=16)
    txt(s, Inches(3.1), y + Inches(0.14), Inches(9.4), Inches(1.05),
        [[(rest, INK, False)]], size=13.5)
    y += Inches(1.45)
txt(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.9),
    [[("A rooftop in Tay Ho charges a Grab driver's battery at noon — morning sunshine "
       "becomes the clean ride home tonight.", GREEN, True)],
     [("⚡ FirmGrid · [TEAM_NAME] · demo: streamlit run app.py", MUTED, False)]], size=16)
notes(s, "WRAP UP (30s). Re-emphasize three things: the impact number, the market fit — "
         "legal today, no new law, pilot path already mapped — and the USP: predict, "
         "allocate fairly, recruit mobility batteries. Land the one-liner: morning sunshine "
         "becomes the clean ride home tonight — and, at scale, firm 24/7 power for the "
         "digital economy. Invite questions; the prototype is open.")

# slide numbers, bottom-right, above the green strip
total = len(prs.slides)
for i, s in enumerate(prs.slides, 1):
    txt(s, W - Inches(1.35), H - Inches(0.45), Inches(1.05), Inches(0.3),
        [[(f"{i} / {total}", MUTED, False)]], size=11, align=PP_ALIGN.RIGHT)

deck_path = OUT / "[TEAM_NAME]_FirmGrid_Slides.pptx"
prs.save(deck_path)
print(f"[deck] {deck_path} ({len(prs.slides)} slides)")

# --------------------------------------------------------------------- #
# poster (single 16:9 slide for the TV display)
# --------------------------------------------------------------------- #
pp = Presentation()
pp.slide_width, pp.slide_height = W, H
s = add_slide(pp)
txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.0),
    [[("⚡ FirmGrid", GREEN, True)]], size=44)
txt(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.9),
    [[("Vietnam throws away clean power at noon and charges its clean vehicles on coal "
       "at night. FirmGrid closes both loops.", INK, True)]], size=19)
y = Inches(2.6)
stat_tile(s, Inches(0.6), y, Inches(2.9), Inches(1.75), "703 → 20",
          "kWh wasted on one real Hanoi feeder-day: baseline vs FirmGrid ON — 97% recovered",
          GREEN, 26)
stat_tile(s, Inches(3.67), y, Inches(2.9), Inches(1.75), "15–19 GWh",
          "recovered per year across 1,000 Hanoi transformers ≈ 10–13 kt CO₂/yr", AMBER, 26)
stat_tile(s, Inches(6.74), y, Inches(2.9), Inches(1.75), "≈1 kg CO₂",
          "avoided per battery swap charged in the solar window", BLUE, 26)
stat_tile(s, Inches(9.81), y, Inches(2.9), Inches(1.75), "54% · $85",
          "hour-matched CFE and $/MWh for a 20 MW data-centre firm block (vs $92 grid)", GREEN, 26)
panel(s, Inches(0.6), Inches(4.6), Inches(7.9), Inches(2.35))
txt(s, Inches(0.85), Inches(4.75), Inches(7.5), Inches(0.45),
    [[("FORECAST → AUCTION → STEER → SETTLE, every 15 minutes", BLUE, True)]], size=14)
bullets(s, Inches(0.85), Inches(5.25), Inches(7.5), Inches(1.6), [
    ("GridMind ", "predicts safe headroom per transformer (F1 0.86 on real Hanoi weather)."),
    ("HeadRoom Auction ", "allocates it fairly — minimal, explainable curtailment."),
    ("FlexMatch ", "pays swap stations to charge on sunshine, not coal."),
    ("TrustLedger ", "verifies, pays VND, seals a tamper-evident audit."),
], size=12.5, gap=4)
panel(s, Inches(8.7), Inches(4.6), Inches(4.0), Inches(2.35), edge=GREEN)
txt(s, Inches(8.95), Inches(4.75), Inches(3.5), Inches(0.5),
    [[("TWO TIERS, ONE ENGINE", GREEN, True)]], size=14)
txt(s, Inches(8.95), Inches(5.25), Inches(3.5), Inches(1.6),
    [[("Tier 1 · Sun-to-Wheels — live today, legal today (Decrees 58 & 57/2025, VND rails).",
       INK, False)],
     [("Tier 2 · Sun-to-Servers — firm 24/7 blocks for data centres via DPPA.", INK, False)],
     [("[TEAM_NAME] · Track 1 · 2026", MUTED, False)]], size=12.5)
poster_path = OUT / "[TEAM_NAME]_FirmGrid_Poster.pptx"
pp.save(poster_path)
print(f"[poster] {poster_path}")
